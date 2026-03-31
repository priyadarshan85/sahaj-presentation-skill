#!/usr/bin/env python3
"""
Sahaj Presentation Generator

Reads a JSON slide specification from stdin and produces a .pptx file
matching Sahaj's corporate visual identity (derived from BusinessCase.pptx).

Usage:
    cat spec.json | python generate_presentation.py
"""

import json
import sys
import os
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree
from lxml.etree import SubElement

# ─── Brand Constants ───────────────────────────────────────────────────────────

SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(5.625)

# Colors
COLOR_DARK_BLUE = RGBColor(0x00, 0x20, 0x60)      # #002060 — titles, section dividers
COLOR_BODY_TEXT = RGBColor(0x0E, 0x0E, 0x0E)       # #0E0E0E — body text
COLOR_CARD_HEADING = RGBColor(0x60, 0x61, 0xAD)    # #6061AD — card headings (purple)
COLOR_ACCENT_BLUE = RGBColor(0x0C, 0x53, 0x95)     # #0C5395 — links/accents
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_GRAY = RGBColor(0x66, 0x66, 0x66)
COLOR_BULLET = RGBColor(0x0E, 0x0E, 0x0E)

# Fonts
FONT_HEADING = "Zilla Slab"
FONT_BODY = "Mulish"

# Layout positions (in inches, derived from BusinessCase.pptx analysis)
TITLE_LEFT = 0.37
TITLE_TOP = 0.15
TITLE_WIDTH = 5.68
TITLE_HEIGHT = 0.61

SUBTITLE_LEFT = 0.37
SUBTITLE_TOP = 0.64
SUBTITLE_WIDTH = 8.54
SUBTITLE_HEIGHT = 0.42

BODY_LEFT = 0.35
BODY_TOP = 1.11
BODY_WIDTH = 8.65
BODY_HEIGHT = 4.03

# Card grid constants
CARD_TOP = 1.11
CARD_LEFT_START = 0.375
CARD_GAP = 0.3
CARD_ROW_GAP = 0.45

# Sahaj logo position (from slide 1 analysis)
SAHAJ_LOGO_LEFT = 3.05
SAHAJ_LOGO_TOP = 2.92
SAHAJ_LOGO_WIDTH = 2.50
SAHAJ_LOGO_HEIGHT = 0.70

CLIENT_LOGO_LEFT = 1.01
CLIENT_LOGO_TOP = 2.97
CLIENT_LOGO_WIDTH = 1.39
CLIENT_LOGO_HEIGHT = 0.62

SEPARATOR_LEFT = 2.88
SEPARATOR_TOP = 2.80
SEPARATOR_HEIGHT = 0.94

SKILL_DIR = Path(__file__).parent.parent
ASSETS_DIR = SKILL_DIR / "assets"


# ─── XML helpers ───────────────────────────────────────────────────────────────

NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}


def _set_font_xml(run_elem, font_name, size_pt, bold=False, color=None):
    """Set font properties via XML for reliable font control."""
    rPr = run_elem.find(f'{{{NSMAP["a"]}}}rPr')
    if rPr is None:
        rPr = SubElement(run_elem, f'{{{NSMAP["a"]}}}rPr')

    rPr.set("lang", "en-US")
    rPr.set("sz", str(int(size_pt * 100)))
    if bold:
        rPr.set("b", "1")

    if color:
        # Remove existing solidFill
        for sf in rPr.findall(f'{{{NSMAP["a"]}}}solidFill'):
            rPr.remove(sf)
        solid_fill = SubElement(rPr, f'{{{NSMAP["a"]}}}solidFill')
        srgb = SubElement(solid_fill, f'{{{NSMAP["a"]}}}srgbClr')
        srgb.set("val", f"{color[0]:02X}{color[1]:02X}{color[2]:02X}")

    # Set latin, ea, cs fonts
    for tag in ["latin", "ea", "cs"]:
        existing = rPr.findall(f'{{{NSMAP["a"]}}}{tag}')
        for e in existing:
            rPr.remove(e)
        font_elem = SubElement(rPr, f'{{{NSMAP["a"]}}}{tag}')
        font_elem.set("typeface", font_name)


def _set_bullet_xml(paragraph, level=0, bullet_char="●", bullet_color=None, font_name=FONT_BODY, size_pt=12):
    """Set bullet properties via XML."""
    pPr = paragraph._p.find(f'{{{NSMAP["a"]}}}pPr')
    if pPr is None:
        pPr = SubElement(paragraph._p, f'{{{NSMAP["a"]}}}pPr')
        paragraph._p.insert(0, pPr)

    pPr.set("lvl", str(level))

    if level == 0:
        pPr.set("indent", str(Emu(Inches(-0.3333)).emu if hasattr(Inches(-0.3333), 'emu') else -304800))
        pPr.set("marL", str(457200))  # 0.5 inches
    elif level == 1:
        pPr.set("indent", str(-298450))
        pPr.set("marL", str(914400))  # 1.0 inches
    else:
        pPr.set("indent", str(-298450))
        pPr.set("marL", str(1371600))  # 1.5 inches

    # Remove existing bullet settings
    for tag in ["buNone", "buChar", "buAutoNum"]:
        for e in pPr.findall(f'{{{NSMAP["a"]}}}{tag}'):
            pPr.remove(e)

    # Set bullet font
    buFont = pPr.find(f'{{{NSMAP["a"]}}}buFont')
    if buFont is None:
        buFont = SubElement(pPr, f'{{{NSMAP["a"]}}}buFont')
    buFont.set("typeface", font_name)

    # Set bullet color
    for bc in pPr.findall(f'{{{NSMAP["a"]}}}buClr'):
        pPr.remove(bc)
    if bullet_color:
        buClr = SubElement(pPr, f'{{{NSMAP["a"]}}}buClr')
        srgb = SubElement(buClr, f'{{{NSMAP["a"]}}}srgbClr')
        srgb.set("val", f"{bullet_color[0]:02X}{bullet_color[1]:02X}{bullet_color[2]:02X}")

    # Set bullet size
    for bs in pPr.findall(f'{{{NSMAP["a"]}}}buSzPts'):
        pPr.remove(bs)
    buSzPts = SubElement(pPr, f'{{{NSMAP["a"]}}}buSzPts')
    buSzPts.set("val", str(int(size_pt * 100)))

    # Set bullet character
    buChar = pPr.find(f'{{{NSMAP["a"]}}}buChar')
    if buChar is None:
        buChar = SubElement(pPr, f'{{{NSMAP["a"]}}}buChar')
    buChar.set("char", bullet_char)

    # Spacing
    for sp_tag in ["spcBef", "spcAft"]:
        for e in pPr.findall(f'{{{NSMAP["a"]}}}{sp_tag}'):
            pPr.remove(e)
    spcBef = SubElement(pPr, f'{{{NSMAP["a"]}}}spcBef')
    spcPts = SubElement(spcBef, f'{{{NSMAP["a"]}}}spcPts')
    spcPts.set("val", "1200")  # 12pt before

    spcAft = SubElement(pPr, f'{{{NSMAP["a"]}}}spcAft')
    spcPts2 = SubElement(spcAft, f'{{{NSMAP["a"]}}}spcPts')
    spcPts2.set("val", "0")

    # Line spacing 115%
    for ls in pPr.findall(f'{{{NSMAP["a"]}}}lnSpc'):
        pPr.remove(ls)
    lnSpc = SubElement(pPr, f'{{{NSMAP["a"]}}}lnSpc')
    spcPct = SubElement(lnSpc, f'{{{NSMAP["a"]}}}spcPct')
    spcPct.set("val", "115000")


def _add_text_run(paragraph, text, font_name, size_pt, bold=False, color=None):
    """Add a text run to a paragraph with full font control."""
    run = paragraph.add_run()
    run.text = text
    _set_font_xml(run._r, font_name, size_pt, bold=bold, color=color)
    return run


def _add_text_with_bold_markers(paragraph, text, font_name, size_pt, color=None):
    """Parse **bold** markers in text and add runs accordingly."""
    parts = re.split(r'(\*\*.*?\*\*)', text)
    for part in parts:
        if part.startswith("**") and part.endswith("**"):
            _add_text_run(paragraph, part[2:-2], font_name, size_pt, bold=True, color=color)
        elif part:
            _add_text_run(paragraph, part, font_name, size_pt, bold=False, color=color)


# ─── Slide Builders ───────────────────────────────────────────────────────────

def add_title_slide(prs, spec):
    """Title slide with logos and presentation title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title text
    title_left = Inches(3.24)
    title_top = Inches(4.15)
    title_width = Inches(6.38)
    title_height = Inches(0.84)

    txBox = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    _add_text_run(p, spec.get("title", "Presentation Title"), FONT_HEADING, 23, bold=True, color=COLOR_DARK_BLUE)

    # Sahaj logo
    sahaj_logo = ASSETS_DIR / "sahaj_logo.png"
    if sahaj_logo.exists():
        slide.shapes.add_picture(
            str(sahaj_logo),
            Inches(SAHAJ_LOGO_LEFT), Inches(SAHAJ_LOGO_TOP),
            Inches(SAHAJ_LOGO_WIDTH), Inches(SAHAJ_LOGO_HEIGHT)
        )

    # Vertical separator line
    from pptx.util import Emu as EmuUtil
    line = slide.shapes.add_connector(
        1,  # straight connector
        Inches(SEPARATOR_LEFT), Inches(SEPARATOR_TOP),
        Inches(SEPARATOR_LEFT), Inches(SEPARATOR_TOP + SEPARATOR_HEIGHT)
    )
    line.line.color.rgb = COLOR_GRAY
    line.line.width = Pt(1)

    # Client logo (if provided)
    client_logo_path = spec.get("client_logo_path")
    if client_logo_path and os.path.exists(client_logo_path):
        slide.shapes.add_picture(
            client_logo_path,
            Inches(CLIENT_LOGO_LEFT), Inches(CLIENT_LOGO_TOP),
            Inches(CLIENT_LOGO_WIDTH), Inches(CLIENT_LOGO_HEIGHT)
        )

    return slide


def add_section_divider(prs, spec):
    """Dark blue section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Set background to dark blue
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_DARK_BLUE

    # Centered title
    txBox = slide.shapes.add_textbox(
        Inches(2.14), Inches(2.47),
        Inches(5.95), Inches(0.69)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _add_text_run(p, spec.get("title", "Section Title"), FONT_HEADING, 27, bold=True, color=COLOR_WHITE)

    return slide


def _add_slide_header(slide, title, subtitle=None):
    """Add the standard title + subtitle header to a content slide."""
    # Title
    txBox = slide.shapes.add_textbox(
        Inches(TITLE_LEFT), Inches(TITLE_TOP),
        Inches(TITLE_WIDTH), Inches(TITLE_HEIGHT)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    _add_text_run(p, title, FONT_HEADING, 24, bold=True, color=COLOR_DARK_BLUE)

    # Subtitle
    if subtitle:
        txBox2 = slide.shapes.add_textbox(
            Inches(SUBTITLE_LEFT), Inches(SUBTITLE_TOP),
            Inches(SUBTITLE_WIDTH), Inches(SUBTITLE_HEIGHT)
        )
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        _add_text_run(p2, subtitle, FONT_HEADING, 14, bold=False, color=COLOR_DARK_BLUE)


def add_bullet_content(prs, spec):
    """Content slide with hierarchical bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    _add_slide_header(slide, spec.get("title", ""), spec.get("subtitle"))

    # Body text box
    txBox = slide.shapes.add_textbox(
        Inches(BODY_LEFT), Inches(BODY_TOP),
        Inches(BODY_WIDTH), Inches(BODY_HEIGHT)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    bullets = spec.get("bullets", [])
    for i, bullet in enumerate(bullets):
        level = bullet.get("level", 0)
        text = bullet.get("text", "")

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Bullet character based on level
        if level == 0:
            char = "●"
        elif level == 1:
            char = "○"
        else:
            char = "-"

        _set_bullet_xml(p, level=level, bullet_char=char, bullet_color=COLOR_BULLET, font_name=FONT_BODY, size_pt=12)
        _add_text_with_bold_markers(p, text, FONT_BODY, 12, color=COLOR_BODY_TEXT)

    return slide


def add_card_content(prs, spec):
    """Content slide with card grid layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    _add_slide_header(slide, spec.get("title", ""), spec.get("subtitle"))

    cards = spec.get("cards", [])
    if not cards:
        return slide

    num_cards = len(cards)

    # Determine grid layout
    if num_cards <= 3:
        cols = num_cards
        rows = 1
    elif num_cards <= 6:
        cols = min(3, (num_cards + 1) // 2)
        rows = (num_cards + cols - 1) // cols
    else:
        cols = 3
        rows = (num_cards + 2) // 3

    # Calculate card dimensions
    available_width = 8.65 - CARD_LEFT_START + 0.375  # total body width
    card_width = (available_width - (cols - 1) * CARD_GAP) / cols

    # Row heights
    if rows == 1:
        card_height = 2.5
    else:
        available_height = 5.0 - CARD_TOP  # leave some bottom margin
        card_height = (available_height - (rows - 1) * CARD_ROW_GAP) / rows

    card_idx = 0
    for row in range(rows):
        cards_in_row = min(cols, num_cards - card_idx)
        for col in range(cards_in_row):
            if card_idx >= num_cards:
                break

            card = cards[card_idx]
            left = CARD_LEFT_START + col * (card_width + CARD_GAP)
            top = CARD_TOP + row * (card_height + CARD_ROW_GAP)

            txBox = slide.shapes.add_textbox(
                Inches(left), Inches(top),
                Inches(card_width), Inches(card_height)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.auto_size = None

            # Card heading
            p_heading = tf.paragraphs[0]
            _add_text_run(p_heading, card.get("heading", ""), FONT_BODY, 11, bold=True, color=COLOR_CARD_HEADING)

            # Card body text
            body_text = card.get("body", "")
            if body_text:
                p_body = tf.add_paragraph()
                _add_text_with_bold_markers(p_body, body_text, FONT_BODY, 11, color=COLOR_BODY_TEXT)

            # Card sub-bullets (optional)
            for sub_bullet in card.get("bullets", []):
                p_sub = tf.add_paragraph()
                _set_bullet_xml(p_sub, level=0, bullet_char="-", bullet_color=COLOR_BODY_TEXT, font_name=FONT_BODY, size_pt=11)
                _add_text_run(p_sub, sub_bullet, FONT_BODY, 11, color=COLOR_BODY_TEXT)

            card_idx += 1

    return slide


# ─── Main ──────────────────────────────────────────────────────────────────────

SLIDE_BUILDERS = {
    "title": add_title_slide,
    "section_divider": add_section_divider,
    "bullet_content": add_bullet_content,
    "card_content": add_card_content,
}


def main():
    spec = json.load(sys.stdin)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    output_path = spec.get("output_path", "presentation.pptx")

    for slide_spec in spec.get("slides", []):
        slide_type = slide_spec.get("type", "bullet_content")
        builder = SLIDE_BUILDERS.get(slide_type)
        if builder:
            builder(prs, slide_spec)
        else:
            print(f"Warning: Unknown slide type '{slide_type}', skipping.", file=sys.stderr)

    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
